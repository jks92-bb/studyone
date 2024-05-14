from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]

title1.text = "실업경험이 가계소비에 미치는 장기효과 분석"
subtitle1.text = "한국의 사례\n작성자: 최영준"

# Slide 2: 연구 배경
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "연구 배경"
content2.text = "• 1997년 외환위기 이후 가계소비 둔화\n• 상흔 소비(scarred consumption) 개념 소개"

# Slide 3: 연구 목적 및 방법
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "연구 목적 및 방법"
content3.text = "• 연구 목적: 실업경험이 가계소비에 미치는 장기적 영향 분석\n• 연구 방법: 미시 패널 자료를 사용한 실증분석"

# Slide 4: 주요 결과
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]

title4.text = "주요 결과"
content4.text = "• 과거 실업경험의 부정적 영향\n• 저축 증가 경로를 통한 상흔 소비 발생"

# Slide 5: 계층별 및 재화별 영향
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
content5 = slide5.placeholders[1]

title5.text = "계층별 및 재화별 영향"
content5.text = "• 저소득 및 저자산 계층에서 두드러짐\n• 선택재와 비내구재 소비에서 뚜렷"

# Slide 6: 결론 및 정책적 시사점
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
content6 = slide6.placeholders[1]

title6.text = "결론 및 정책적 시사점"
content6.text = "• 상흔 소비의 존재 확인\n• 취약계층 지원의 필요성 강조"

# Save the presentation
prs.save("실업경험이_가계소비에_미치는_장기효과_분석.pptx")
